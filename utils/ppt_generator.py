"""
PulseAI - Automated RBI-Style PowerPoint Generator
One-click boardroom presentations with premium styling
"""

import io
import os
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio


# RBI Brand Colors
RBI_NAVY = RGBColor(15, 27, 61)      # #0f1b3d
RBI_GOLD = RGBColor(212, 175, 55)    # #d4af37
RBI_WHITE = RGBColor(255, 255, 255)
RBI_LIGHT_BLUE = RGBColor(66, 103, 178)
RBI_GRAY = RGBColor(128, 128, 128)


class PulseAIPresentationGenerator:
    def __init__(self):
        """Initialize presentation with RBI-style template"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Define layouts
        self.blank_layout = self.prs.slide_layouts[6]  # Blank
        self.title_layout = self.prs.slide_layouts[0]  # Title
    
    def add_title_slide(self, month_year):
        """Create title slide"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        
        # Navy background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RBI_NAVY
        
        # PulseAI Logo placeholder (top-left)
        logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.6))
        logo_frame = logo_box.text_frame
        logo_frame.text = "PulseAI"
        p = logo_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RBI_GOLD
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = f"Indian Financial Intelligence Report"
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RBI_WHITE
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = month_year
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.color.rgb = RBI_GOLD
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.4))
        footer_frame = footer_box.text_frame
        footer_frame.text = "Powered by PulseAI | Data: RBI, NPCI, NSE, AMFI"
        p = footer_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = RBI_GRAY
    
    def add_executive_summary_slide(self, summary_text):
        """Add executive summary slide"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(250, 250, 255)
        
        # Header bar
        header = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), Inches(10), Inches(0.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RBI_NAVY
        header.line.fill.background()
        
        # Header text
        header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        frame = header_text.text_frame
        frame.text = "Executive Summary"
        p = frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RBI_GOLD
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = summary_text
        
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(30, 30, 30)
            paragraph.space_after = Pt(12)
    
    def add_data_slide(self, title, content_items):
        """Add slide with bullet points"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(250, 250, 255)
        
        # Header
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = RBI_NAVY
        header.line.fill.background()
        
        header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        frame = header_text.text_frame
        frame.text = title
        p = frame.paragraphs[0]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RBI_GOLD
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for item in content_items:
            p = content_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(30, 30, 30)
            p.space_after = Pt(14)
    
    def add_chart_slide(self, title, chart_image_bytes):
        """Add slide with chart"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(250, 250, 255)
        
        # Header
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = RBI_NAVY
        header.line.fill.background()
        
        header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        frame = header_text.text_frame
        frame.text = title
        p = frame.paragraphs[0]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RBI_GOLD
        
        # Chart image
        if chart_image_bytes:
            slide.shapes.add_picture(
                chart_image_bytes,
                Inches(1), Inches(1.5), 
                width=Inches(8), height=Inches(5)
            )
    
    def add_closing_slide(self):
        """Add final branded slide"""
        slide = self.prs.slides.add_slide(self.blank_layout)
        
        # Gold background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RBI_NAVY
        
        # Main message
        msg_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        msg_frame = msg_box.text_frame
        msg_frame.text = "Thank You"
        p = msg_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RBI_GOLD
        
        # Branding
        brand_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
        brand_frame = brand_box.text_frame
        brand_frame.text = "PulseAI\nReal-Time Indian Financial Intelligence"
        
        for para in brand_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(200, 200, 200)
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.4))
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("Generated on %B %d, %Y")
        p = date_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = RBI_GRAY
    
    def save(self, filename=None):
        """Save presentation"""
        if filename is None:
            filename = f"PulseAI_Report_{datetime.now().strftime('%Y%m%d')}.pptx"
        
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output, filename


def generate_boardroom_presentation(data_dict, executive_summary, anomalies, forecasts):
    """
    Generate complete boardroom presentation
    
    Args:
        data_dict: Dictionary with all financial datasets
        executive_summary: AI-generated summary text
        anomalies: List of detected anomalies
        forecasts: Dictionary with forecast data
    
    Returns:
        BytesIO object with PPT, filename
    """
    gen = PulseAIPresentationGenerator()
    
    # Slide 1: Title
    month_year = datetime.now().strftime("%B %Y")
    gen.add_title_slide(month_year)
    
    # Slide 2-3: Executive Summary
    gen.add_executive_summary_slide(executive_summary)
    
    # Slide 4: UPI Highlights
    if 'upi' in data_dict and not data_dict['upi'].empty:
        upi_data = data_dict['upi'].tail(1).iloc[0]
        upi_content = [
            f"✓ Transaction Volume: {upi_data['Volume_Billion']:.2f} billion transactions",
            f"✓ Transaction Value: ₹{upi_data['Value_LakhCrore']:.2f} lakh crore",
            f"✓ Average Transaction: ₹{upi_data['Avg_Transaction_Size']:.2f}",
            f"✓ YoY Growth: Strong upward trajectory",
            f"✓ Peak adoption during festival seasons"
        ]
        gen.add_data_slide("UPI Transaction Highlights", upi_content)
    
    # Slide 5: State-wise Banking
    if 'rbi_credit' in data_dict and not data_dict['rbi_credit'].empty:
        top_states = data_dict['rbi_credit'].nlargest(5, 'Credit_Growth_%')
        banking_content = [
            f"Top Growing States:",
            *[f"• {row['State']}: {row['Credit_Growth_%']:.1f}% growth, CD Ratio: {row['CD_Ratio']:.1f}%" 
              for _, row in top_states.iterrows()],
            "",
            f"Average National Credit Growth: {data_dict['rbi_credit']['Credit_Growth_%'].mean():.1f}%",
            f"Highest Digital Adoption: {data_dict['rbi_credit'].nlargest(1, 'Digital_Adoption_%').iloc[0]['State']}"
        ]
        gen.add_data_slide("State-wise Banking Performance", banking_content)
    
    # Slide 6: Stock Market Snapshot
    if 'nse' in data_dict and not data_dict['nse'].empty:
        nse_content = [
            f"Top Performing Stocks:",
            *[f"• {row['Symbol']}: ₹{row['LTP']:.2f} ({row['Change_%']:+.2f}%)" 
              for _, row in data_dict['nse'].head(5).iterrows()],
            "",
            f"Market Sentiment: {'Positive' if data_dict['nse']['Change_%'].mean() > 0 else 'Cautious'}",
            f"Average Change: {data_dict['nse']['Change_%'].mean():.2f}%"
        ]
        gen.add_data_slide("NSE Market Snapshot", nse_content)
    
    # Slide 7: Mutual Funds
    if 'mutual_funds' in data_dict and not data_dict['mutual_funds'].empty:
        latest = data_dict['mutual_funds'][data_dict['mutual_funds']['Month'] == data_dict['mutual_funds']['Month'].max()]
        mf_content = [
            f"Industry AUM: ₹{latest['AUM_LakhCrore'].sum():.2f} lakh crore",
            "",
            "Category-wise Breakdown:",
            *[f"• {row['Category']}: ₹{row['AUM_LakhCrore']:.2f} lakh crore" 
              for _, row in latest.nlargest(5, 'AUM_LakhCrore').iterrows()],
            "",
            f"Total Investor Accounts: {latest['Accounts_Lakh'].sum():.2f} lakh"
        ]
        gen.add_data_slide("Mutual Fund Industry", mf_content)
    
    # Slide 8: RBI Policy
    if 'rbi_policy' in data_dict and not data_dict['rbi_policy'].empty:
        latest_policy = data_dict['rbi_policy'].iloc[-1]
        policy_content = [
            f"Current Repo Rate: {latest_policy['Repo_Rate']}%",
            f"Reverse Repo Rate: {latest_policy['Reverse_Repo']}%",
            f"CRR: {latest_policy['CRR']}%",
            f"SLR: {latest_policy['SLR']}%",
            "",
            f"Policy Stance: {latest_policy['Policy_Stance']}",
            "",
            "Implication: Stable rates supporting growth while managing inflation"
        ]
        gen.add_data_slide("RBI Monetary Policy", policy_content)
    
    # Slide 9: Anomalies
    anomaly_content = anomalies if anomalies else ["No significant anomalies detected in current period"]
    gen.add_data_slide("Anomalies & Alerts", anomaly_content)
    
    # Slide 10: Forecasts
    if forecasts:
        forecast_content = []
        for metric, data in forecasts.items():
            forecast_content.append(f"{metric}: {data.get('narrative', 'Steady growth expected')[:100]}...")
        gen.add_data_slide("30-Day Forecasts", forecast_content)
    
    # Slide 11: Key Takeaways
    takeaways = [
        "✓ UPI continues to be the backbone of digital payments",
        "✓ Regional banking growth shows healthy diversification",
        "✓ Stock markets remain resilient despite global headwinds",
        "✓ Mutual fund industry witnessing strong retail participation",
        "✓ RBI's accommodative stance supporting economic recovery"
    ]
    gen.add_data_slide("Key Takeaways", takeaways)
    
    # Slide 12: Closing
    gen.add_closing_slide()
    
    # Save and return
    return gen.save()


if __name__ == "__main__":
    # Test presentation generation
    print("Testing PPT Generator...")
    gen = PulseAIPresentationGenerator()
    gen.add_title_slide("November 2025")
    gen.add_executive_summary_slide("This is a test summary.")
    gen.add_data_slide("Test Slide", ["Point 1", "Point 2", "Point 3"])
    gen.add_closing_slide()
    output, filename = gen.save("test_presentation.pptx")
    print(f"✅ Generated {filename}")
